#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Create Comprehensive PowerPoint Presentation for SAE Gender Bias VLM Project
=============================================================================

This script generates a detailed presentation covering all aspects of the research project.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
import os

# Helper function for RGB colors
def RgbColor(r, g, b):
    """Create RGB color compatible with python-pptx."""
    from pptx.dml.color import RGBColor
    return RGBColor(r, g, b)

# Color scheme
COLORS = {
    'title_bg': RgbColor(30, 58, 138),      # Dark blue
    'accent': RgbColor(59, 130, 246),       # Blue
    'paligemma': RgbColor(234, 88, 12),     # Orange
    'qwen': RgbColor(22, 163, 74),          # Green
    'llama': RgbColor(147, 51, 234),        # Purple
    'dark': RgbColor(31, 41, 55),           # Dark gray
    'light': RgbColor(249, 250, 251),       # Light gray
    'white': RgbColor(255, 255, 255),
    'red': RgbColor(220, 38, 38),
    'green': RgbColor(34, 197, 94),
}

def set_slide_background(slide, color):
    """Set slide background color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_slide(prs, title, subtitle=""):
    """Add a title slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLORS['title_bg'])
    
    # Title
    left = Inches(0.5)
    top = Inches(2.5)
    width = Inches(9)
    height = Inches(1.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    if subtitle:
        top = Inches(4.2)
        height = Inches(1)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = RgbColor(200, 200, 200)
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_section_slide(prs, title, number=""):
    """Add a section divider slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLORS['accent'])
    
    # Section number
    if number:
        left = Inches(0.5)
        top = Inches(2)
        width = Inches(9)
        height = Inches(0.8)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = f"Section {number}"
        p.font.size = Pt(20)
        p.font.color.rgb = RgbColor(200, 220, 255)
        p.alignment = PP_ALIGN.CENTER
    
    # Title
    left = Inches(0.5)
    top = Inches(2.8)
    width = Inches(9)
    height = Inches(1.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, bullets, notes=""):
    """Add a content slide with bullets."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLORS['light'])
    
    # Title bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS['title_bg']
    shape.line.fill.background()
    
    # Title text
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.7)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    
    # Bullets
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        # Handle nested bullets
        if isinstance(bullet, tuple):
            text, level = bullet
            p.text = text
            p.level = level
        else:
            p.text = f"• {bullet}"
            p.level = 0
        
        p.font.size = Pt(20)
        p.font.color.rgb = COLORS['dark']
        p.space_after = Pt(12)
    
    # Speaker notes
    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes
    
    return slide

def add_two_column_slide(prs, title, left_title, left_bullets, right_title, right_bullets):
    """Add a two-column content slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLORS['light'])
    
    # Title bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS['title_bg']
    shape.line.fill.background()
    
    # Title text
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    
    # Left column title
    txBox = slide.shapes.add_textbox(Inches(0.3), Inches(1.4), Inches(4.5), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLORS['accent']
    
    # Left column bullets
    txBox = slide.shapes.add_textbox(Inches(0.3), Inches(1.9), Inches(4.5), Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(left_bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS['dark']
        p.space_after = Pt(8)
    
    # Right column title
    txBox = slide.shapes.add_textbox(Inches(5.2), Inches(1.4), Inches(4.5), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLORS['accent']
    
    # Right column bullets
    txBox = slide.shapes.add_textbox(Inches(5.2), Inches(1.9), Inches(4.5), Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(right_bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS['dark']
        p.space_after = Pt(8)
    
    return slide

def add_table_slide(prs, title, headers, rows):
    """Add a slide with a table."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLORS['light'])
    
    # Title bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS['title_bg']
    shape.line.fill.background()
    
    # Title text
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    
    # Table
    n_cols = len(headers)
    n_rows = len(rows) + 1  # +1 for header
    
    left = Inches(0.3)
    top = Inches(1.5)
    width = Inches(9.4)
    height = Inches(0.5 * n_rows)
    
    table = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table
    
    # Set column widths
    col_width = Inches(9.4 / n_cols)
    for i in range(n_cols):
        table.columns[i].width = col_width
    
    # Header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS['title_bg']
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.CENTER
    
    # Data rows
    for row_idx, row in enumerate(rows):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(value)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.font.color.rgb = COLORS['dark']
            p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_key_finding_slide(prs, finding_number, title, description, evidence):
    """Add a key finding highlight slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLORS['light'])
    
    # Finding number badge
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.3), Inches(0.3), Inches(1), Inches(1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS['accent']
    shape.line.fill.background()
    
    # Number text
    txBox = slide.shapes.add_textbox(Inches(0.3), Inches(0.5), Inches(1), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = str(finding_number)
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    
    # Title
    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(0.4), Inches(8), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"Key Finding #{finding_number}"
    p.font.size = Pt(18)
    p.font.color.rgb = COLORS['accent']
    p = tf.add_paragraph()
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark']
    
    # Description box
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(1.8), Inches(9.4), Inches(2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RgbColor(239, 246, 255)  # Light blue
    shape.line.color.rgb = COLORS['accent']
    
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.6))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = description
    p.font.size = Pt(18)
    p.font.color.rgb = COLORS['dark']
    
    # Evidence
    txBox = slide.shapes.add_textbox(Inches(0.3), Inches(4.2), Inches(9.4), Inches(2.3))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Evidence:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLORS['accent']
    
    for item in evidence:
        p = tf.add_paragraph()
        p.text = f"  {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS['dark']
        p.space_after = Pt(6)
    
    return slide

def create_presentation():
    """Create the full presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # =========================================================================
    # SLIDE 1: Title Slide
    # =========================================================================
    add_title_slide(
        prs,
        "Sparse Autoencoders Reveal and Control\nGender Bias in Vision-Language Models",
        "A Cross-Lingual Mechanistic Interpretability Study\n\nNour Mubarak | Durham University | March 2026"
    )
    
    # =========================================================================
    # SLIDE 2: Project Overview
    # =========================================================================
    add_content_slide(prs, "Project Overview", [
        "First application of Sparse Autoencoders (SAEs) to VLMs for mechanistic interpretability",
        "Goal: Understand HOW gender bias is encoded, not just that it exists",
        "Cross-lingual analysis: English and Arabic",
        "4 state-of-the-art VLM architectures studied",
        "Causal intervention experiments to validate findings",
        "Direct implications for bias mitigation strategies"
    ])
    
    # =========================================================================
    # SLIDE 3: Research Questions
    # =========================================================================
    add_content_slide(prs, "Research Questions", [
        "RQ1: Can SAEs identify interpretable gender-associated features in VLMs?",
        "RQ2: Do these features have CAUSAL influence on gendered output?",
        "RQ3: How does gender encoding differ across model architectures?",
        "RQ4: Are gender features shared across languages (English/Arabic)?",
        "RQ5: What are the implications for bias mitigation?"
    ])
    
    # =========================================================================
    # SECTION: Background
    # =========================================================================
    add_section_slide(prs, "Background & Motivation", "1")
    
    # =========================================================================
    # SLIDE 4: The Problem - Gender Bias in VLMs
    # =========================================================================
    add_content_slide(prs, "The Problem: Gender Bias in VLMs", [
        "Vision-Language Models exhibit systematic gender bias in captioning",
        "Example: Man on skateboard -> 'man performing tricks'",
        "Example: Woman on skateboard -> 'person standing on board'",
        "Models systematically deploy gendered language for some groups",
        "Prior work measures WHAT but not WHY or HOW",
        "Need to understand internal mechanisms for effective mitigation"
    ])
    
    # =========================================================================
    # SLIDE 5: Why Mechanistic Interpretability?
    # =========================================================================
    add_two_column_slide(prs, "Why Mechanistic Interpretability?",
        "Prior Approaches (Black Box)",
        [
            "Measure bias in outputs",
            "Data augmentation",
            "Constrained decoding",
            "Treat model as black box",
            "Cannot explain WHY bias occurs"
        ],
        "Our Approach (Mechanistic)",
        [
            "Open the black box",
            "Identify specific features",
            "Test causal hypotheses",
            "Understand internal encoding",
            "Enable surgical interventions"
        ]
    )
    
    # =========================================================================
    # SLIDE 6: What are Sparse Autoencoders?
    # =========================================================================
    add_content_slide(prs, "What are Sparse Autoencoders (SAEs)?", [
        "Neural network that decomposes dense activations into sparse, interpretable features",
        "Architecture: Encoder -> Sparse hidden -> Decoder",
        "Equation: h = ReLU(W_enc * x + b_enc)",
        "Reconstruction: x_hat = W_dec * h + b_dec",
        "Loss: MSE(x, x_hat) + lambda*||h||_1 (reconstruction + sparsity)",
        "Each feature corresponds to a semantically meaningful concept",
        "Previously applied to text LLMs; we extend to multimodal VLMs"
    ])
    
    # =========================================================================
    # SECTION: Methodology
    # =========================================================================
    add_section_slide(prs, "Methodology", "2")
    
    # =========================================================================
    # SLIDE 7: Models Studied
    # =========================================================================
    add_table_slide(prs, "Vision-Language Models Studied",
        ["Model", "Parameters", "d_model", "Layers", "SAE Features", "Type"],
        [
            ["PaLiGemma-3B", "2.9B", "2,048", "18", "16,384", "Captioning"],
            ["Qwen2-VL-7B", "7.6B", "3,584", "28", "28,672", "Instruction-tuned"],
            ["LLaVA-1.5-7B", "7.1B", "4,096", "32", "32,768", "Instruction-tuned"],
            ["Llama-3.2-11B", "10.6B", "4,096", "40", "32,768", "Instruction-tuned"]
        ]
    )
    
    # =========================================================================
    # SLIDE 8: Dataset
    # =========================================================================
    add_content_slide(prs, "Dataset: Flickr8K Bilingual", [
        "Flickr8K dataset (Hodosh et al., 2013)",
        "8,092 images with bilingual English-Arabic caption pairs",
        "Gender labels derived from caption content (keyword matching)",
        "Male samples: 5,562 | Female samples: 2,047 (ratio 2.7:1)",
        "16 binary gender terms tracked (he, she, him, her, man, woman, etc.)",
        "18 gender-neutral terms tracked for comparison",
        "Dataset reflects real-world male bias common in image datasets"
    ])
    
    # =========================================================================
    # SLIDE 9: Pipeline Overview
    # =========================================================================
    add_content_slide(prs, "Four-Stage Pipeline", [
        "Stage 1: ACTIVATION EXTRACTION",
        "   -> Extract decoder layer activations using PyTorch hooks",
        "   -> Mean-pool across sequence positions",
        "Stage 2: SAE TRAINING",
        "   -> Train sparse autoencoder (8x expansion factor)",
        "   -> Vanilla ReLU + L1 sparsity (lambda = 10^-4)",
        "Stage 3: GENDER FEATURE IDENTIFICATION",
        "   -> Differential activation analysis (male vs female)",
        "   -> Select top-k features (k=100, 0.6% of total)",
        "Stage 4: CAUSAL INTERVENTION",
        "   -> Ablate identified features during generation",
        "   -> Compare to random ablation controls (25 runs)"
    ])
    
    # =========================================================================
    # SLIDE 10: SAE Architecture Details
    # =========================================================================
    add_content_slide(prs, "SAE Architecture & Training", [
        "Expansion factor: 8x (e.g., 2048 -> 16,384 features)",
        "Activation: ReLU (vanilla baseline architecture)",
        "Loss: MSE reconstruction + L1 sparsity penalty",
        "Optimizer: AdamW (lr = 10^-4)",
        "Batch size: 256",
        "Training: 50 epochs with early stopping (patience=10)",
        "Split: 90% train / 10% validation",
        "Note: TopK and JumpReLU are alternatives for future work"
    ])
    
    # =========================================================================
    # SLIDE 11: Intervention Design
    # =========================================================================
    add_content_slide(prs, "Causal Intervention Design", [
        "500 test images (seed=42 for reproducibility)",
        "BASELINE: Generate captions with no modification",
        "TARGETED: Zero contribution of top-100 gender features",
        "RANDOM CONTROL: 25 runs with randomly selected features",
        "Two hook variants used:",
        "   -> Full reconstruction (PaLiGemma): Replace with SAE output",
        "   -> Residual ablation (Qwen, Llama): Subtract feature contribution",
        "Deterministic generation: greedy decoding, no sampling"
    ])
    
    # =========================================================================
    # SECTION: Results
    # =========================================================================
    add_section_slide(prs, "Results", "3")
    
    # =========================================================================
    # SLIDE 12: Main Results Table
    # =========================================================================
    add_table_slide(prs, "Main Intervention Results",
        ["Model", "Baseline", "Targeted", "Change", "Random (mean+/-std)", "Ratio"],
        [
            ["PaLiGemma (L9)", "1,522", "1,277", "-16.1%", "-8.7% +/- 3.9%", "1.8x"],
            ["Qwen2-VL (L12)", "1,315", "1,367", "+3.95%", "-0.56% +/- 1.14%", "7.1x"],
            ["Llama (L20)", "1,355", "1,423", "+5.02%", "-0.84% +/- 0.76%", "6.0x"]
        ]
    )
    
    # =========================================================================
    # SLIDE 13: Key Finding 1
    # =========================================================================
    add_key_finding_slide(prs, 1,
        "Targeted Ablation Works",
        "Ablating just 0.6% of SAE features (100 out of 16,384) produces statistically significant changes in gender-marked language across ALL three models tested. Effect ratios range from 1.8x to 7.1x compared to random controls.",
        [
            "Bootstrap 95% CIs exclude zero for all models",
            "PaLiGemma: CI = [-0.274, -0.174]",
            "Qwen2-VL: CI = [+0.038, +0.200]",
            "Llama: CI = [+0.078, +0.242]",
            "25 random control runs per model for statistical rigor"
        ]
    )
    
    # =========================================================================
    # SLIDE 14: Key Finding 2
    # =========================================================================
    add_key_finding_slide(prs, 2,
        "Direction Divergence: Excitatory vs Inhibitory",
        "The DIRECTION of the effect depends on model architecture. PaLiGemma (captioning-focused) shows DECREASE in gender terms (-16.1%), while Qwen2-VL (+3.95%) and Llama (+5.02%) show INCREASES. This reveals fundamentally different encoding mechanisms.",
        [
            "PaLiGemma: Gender features PRODUCE gendered output -> ablation removes it",
            "Qwen/Llama: Gender features SUPPRESS/REGULATE -> ablation releases it",
            "Analogous to excitatory vs inhibitory circuits in neuroscience",
            "Critical implication: Same intervention has OPPOSITE effects!"
        ]
    )
    
    # =========================================================================
    # SLIDE 15: Excitatory vs Inhibitory Mechanisms
    # =========================================================================
    add_two_column_slide(prs, "Two Mechanisms of Gender Encoding",
        "Excitatory (PaLiGemma)",
        [
            "Smaller, captioning-focused model",
            "Gender features directly produce output",
            "Ablation REMOVES gendered language",
            "Effect: -16.1% gender terms",
            "Like turning off a light switch"
        ],
        "Inhibitory (Qwen, Llama)",
        [
            "Larger, instruction-tuned models",
            "Features regulate/suppress gender",
            "Ablation RELEASES gendered language",
            "Effect: +4-5% gender terms",
            "Like removing a governor/brake"
        ]
    )
    
    # =========================================================================
    # SLIDE 16: Layer Specificity
    # =========================================================================
    add_content_slide(prs, "Layer Specificity (PaLiGemma)", [
        "Gender features are LOCALIZED, not distributed uniformly",
        "L9 (mid-layer, ~50%): Carries the ENTIRE effect",
        "   -> Targeted: -16.1%, Random: -8.7% +/- 3.9%, p = 1.65x10^-21",
        "L17 (penultimate, negative control): NO significant effect",
        "   -> Targeted: -0.8%, Random: -4.1% +/- 5.0%, p = 0.999",
        "L9 + L17 combined: Same as L9 alone (no additional benefit)",
        "Conclusion: Middle layers encode semantic features like gender",
        "Justification: Consistent with prior SAE literature on text LLMs"
    ])
    
    # =========================================================================
    # SLIDE 17: Per-Term Analysis
    # =========================================================================
    add_table_slide(prs, "Per-Term Changes Under Ablation",
        ["Term", "PaLiGemma", "Qwen2-VL", "Llama"],
        [
            ["her", "-73.2%", "+4.1%", "+11.9%"],
            ["him", "-83.3%", "+90.9%", "+5.3%"],
            ["he", "-22.6%", "+0.2%", "-0.2%"],
            ["man", "+10.4%", "+53.6%", "+4.8%"],
            ["woman", "+9.3%", "+43.6%", "-4.3%"],
            ["his", "-0.4%", "+6.3%", "+66.3%"]
        ]
    )
    
    # =========================================================================
    # SLIDE 18: Per-Term Insight
    # =========================================================================
    add_content_slide(prs, "Per-Term Analysis: Key Insight", [
        "Ablation RESHUFFLES which gendered terms appear",
        "It does NOT uniformly suppress all gender terms",
        "PaLiGemma pattern:",
        "   -> Pronouns dramatically reduced (her -73%, him -83%)",
        "   -> Explicit nouns slightly increased (man +10%, woman +9%)",
        "   -> Suggests features encode PRONOUN USAGE, not pure gender",
        "Qwen/Llama pattern:",
        "   -> Nouns increased (man +54%, him +91%)",
        "   -> Features appear to SUPPRESS explicit gendering",
        "Implication: Gender encoding is more complex than simple on/off"
    ])
    
    # =========================================================================
    # SLIDE 19: SAE Quality Metrics
    # =========================================================================
    add_table_slide(prs, "SAE Quality Across Models",
        ["Model", "Cosine Sim", "Explained Var.", "L0 (Active)", "Dead %"],
        [
            ["PaLiGemma", "0.9999", "99.8%", "7,992", "51.2%"],
            ["Qwen2-VL", "0.9965", "66.4%", "2,049", "71.6%"],
            ["Llama", "0.9956", "36.6%", "344", "98.6%"]
        ]
    )
    
    # =========================================================================
    # SLIDE 20: SAE Quality Insight
    # =========================================================================
    add_content_slide(prs, "Robustness to SAE Quality", [
        "SAE quality varies SUBSTANTIALLY across models",
        "PaLiGemma: Excellent (99.8% explained variance)",
        "Qwen2-VL: Good (66.4% explained variance)",
        "Llama: Poor (36.6% explained variance, 98.6% dead features)",
        "KEY FINDING: Intervention works even with low-quality SAEs!",
        "   -> Llama achieves 6.0x effect ratio despite poor SAE",
        "   -> Suggests approach is CONSERVATIVE",
        "   -> Better SAEs would likely yield LARGER effects",
        "Justification: This is a strength - results are robust to SAE training"
    ])
    
    # =========================================================================
    # SECTION: Cross-Lingual Analysis
    # =========================================================================
    add_section_slide(prs, "Cross-Lingual Analysis", "4")
    
    # =========================================================================
    # SLIDE 21: Cross-Lingual Research Question
    # =========================================================================
    add_content_slide(prs, "Cross-Lingual Research Question", [
        "Do English and Arabic share gender bias features?",
        "If YES: Debiasing in English would automatically help Arabic",
        "If NO: Language-specific interventions are required",
        "Our approach:",
        "   -> Compute Jaccard overlap of top-100 gender features",
        "   -> Causal 2x2 ablation: Feature language x Caption language",
        "Languages differ fundamentally:",
        "   -> English: Grammatical gender is minimal",
        "   -> Arabic: Extensive grammatical gender marking"
    ])
    
    # =========================================================================
    # SLIDE 22: Feature Overlap Results
    # =========================================================================
    add_table_slide(prs, "Cross-Lingual Feature Overlap (Jaccard Index)",
        ["Model", "Overlap (k=100)", "Jaccard", "Interpretation"],
        [
            ["PaLiGemma", "57/100", "0.399", "Substantial overlap"],
            ["Qwen2-VL", "0/100", "0.000", "Zero overlap"],
            ["LLaVA", "2/100", "0.010", "Near-zero overlap"],
            ["Llama", "1/100", "0.005", "Near-zero overlap"]
        ]
    )
    
    # =========================================================================
    # SLIDE 23: Key Finding 3
    # =========================================================================
    add_key_finding_slide(prs, 3,
        "Language-Specific Gender Encoding",
        "Cross-lingual feature overlap is NEAR-ZERO (Jaccard < 3%) for most models. Causal 2x2 ablation confirms: ablating English features has NO effect on Arabic captions when overlap is zero, but DOES transfer when overlap is high.",
        [
            "Qwen2-VL: 0% overlap -> NO cross-lingual transfer",
            "PaLiGemma: 57% overlap -> Transfer DOES occur",
            "Feature overlap PREDICTS transferability",
            "Implication: Debiasing in English does NOT transfer to Arabic",
            "Multilingual fairness requires language-specific interventions"
        ]
    )
    
    # =========================================================================
    # SLIDE 24: Causal Cross-Lingual Ablation
    # =========================================================================
    add_table_slide(prs, "Cross-Lingual Causal Ablation (Qwen2-VL)",
        ["Condition", "Features", "Caption", "Delta%", "Significant?"],
        [
            ["EN->EN (within)", "English", "English", "+4.0%", "YES"],
            ["EN->AR (cross)", "English", "Arabic", "-2.5%", "NO"],
            ["AR->EN (cross)", "Arabic", "English", "+0.4%", "NO"],
            ["AR->AR (within)", "Arabic", "Arabic", "-6.4%", "NO"]
        ]
    )
    
    # =========================================================================
    # SLIDE 25: Cross-Lingual Insight
    # =========================================================================
    add_content_slide(prs, "Cross-Lingual Causal Confirmation", [
        "Qwen2-VL (0% feature overlap, Jaccard = 0.000):",
        "   -> Only EN->EN is significant",
        "   -> NO cross-lingual transfer in either direction",
        "   -> CAUSAL confirmation of language-specific encoding",
        "PaLiGemma (57% feature overlap, Jaccard = 0.399):",
        "   -> AR->EN IS significant (transfers!)",
        "   -> High overlap predicts cross-lingual transfer",
        "Conclusion: Feature overlap PREDICTS transferability",
        "   -> Zero overlap -> No transfer",
        "   -> High overlap -> Transfer occurs"
    ])
    
    # =========================================================================
    # SECTION: Implications & Discussion
    # =========================================================================
    add_section_slide(prs, "Implications & Discussion", "5")
    
    # =========================================================================
    # SLIDE 26: Practical Implications
    # =========================================================================
    add_content_slide(prs, "Practical Implications for Bias Mitigation", [
        "CRITICAL: Same intervention has OPPOSITE effects across architectures!",
        "Naive debiasing (ablate top-differential features):",
        "   -> REDUCES bias in PaLiGemma (checkmark)",
        "   -> INCREASES bias in Qwen/Llama (warning)",
        "Architecture-specific analysis is ESSENTIAL before deployment",
        "Language-specific interventions needed for multilingual fairness",
        "Cannot assume cross-lingual transfer without measuring overlap",
        "SAE-based analysis should precede any mechanistic intervention"
    ])
    
    # =========================================================================
    # SLIDE 27: Why the Direction Difference?
    # =========================================================================
    add_content_slide(prs, "Hypothesis: Why Direction Divergence?", [
        "Captioning models (PaLiGemma):",
        "   -> Trained to describe images directly",
        "   -> Learn features that PRODUCE gendered descriptions",
        "   -> Gender features are excitatory",
        "Instruction-tuned models (Qwen, Llama):",
        "   -> Trained with RLHF/alignment",
        "   -> Learn features that REGULATE/SUPPRESS gender",
        "   -> May have developed 'safety circuits'",
        "   -> Gender features are inhibitory",
        "Supporting evidence: Both instruction-tuned models show same direction",
        "Alternative: Could partially reflect different hook methods"
    ])
    
    # =========================================================================
    # SLIDE 28: Limitations
    # =========================================================================
    add_content_slide(prs, "Limitations", [
        "SAE training scale: 5K-10K samples (larger sets may help)",
        "SAE architecture: Vanilla ReLU only (TopK/JumpReLU not tested)",
        "Binary gender: Non-binary terms tracked but not primary metric",
        "Hook inconsistency: Different methods for different models",
        "No human evaluation: Caption fluency not assessed",
        "Feature polysemanticity: Features may encode multiple concepts",
        "Llama SAE quality: Only 36.6% explained variance",
        "No SAEBench evaluation: Standardized metrics not reported"
    ])
    
    # =========================================================================
    # SLIDE 29: Future Work
    # =========================================================================
    add_content_slide(prs, "Future Work", [
        "Train larger-scale SAEs (100K+ samples)",
        "Compare SAE architectures (TopK, JumpReLU)",
        "Extend to more languages beyond English/Arabic",
        "Human evaluation of caption quality after ablation",
        "Apply to other bias types (race, age, profession)",
        "Develop architecture-aware debiasing strategies",
        "Investigate instruction tuning's role in creating regulatory features",
        "SAEBench evaluation for standardized comparisons"
    ])
    
    # =========================================================================
    # SECTION: Conclusions
    # =========================================================================
    add_section_slide(prs, "Conclusions", "6")
    
    # =========================================================================
    # SLIDE 30: Summary of Contributions
    # =========================================================================
    add_content_slide(prs, "Summary of Contributions", [
        "METHODOLOGICAL: First SAE-based mechanistic interpretability of VLM bias",
        "EMPIRICAL: Causal evidence across 3 architectures",
        "   -> 0.6% of features produce 1.8-7.1x effects",
        "MECHANISTIC: Discovery of excitatory vs inhibitory encoding",
        "   -> Different architectures use fundamentally different mechanisms",
        "CROSS-LINGUAL: First causal 2x2 ablation on VLM gender features",
        "   -> Feature overlap predicts cross-lingual transferability",
        "PRACTICAL: Architecture & language-specific analysis essential",
        "   -> Same intervention can increase or decrease bias"
    ])
    
    # =========================================================================
    # SLIDE 31: Key Takeaways
    # =========================================================================
    add_content_slide(prs, "Key Takeaways", [
        "1. SAEs can identify causal gender features in VLMs",
        "2. Direction of effect depends on architecture",
        "   -> Captioning models: excitatory features",
        "   -> Instruction-tuned: inhibitory features", 
        "3. Cross-lingual features are language-specific",
        "   -> Zero overlap -> no transfer",
        "   -> High overlap -> transfer occurs",
        "4. Intervention is robust to SAE quality",
        "5. Bias mitigation requires architecture-specific strategies"
    ])
    
    # =========================================================================
    # SLIDE 32: Thank You
    # =========================================================================
    slide = add_title_slide(
        prs,
        "Thank You",
        "Questions?\n\nCode: github.com/nour-mubarak/mechanistic_intrep\nnour.mubarak@durham.ac.uk"
    )
    
    # =========================================================================
    # SUPPLEMENTARY SLIDES
    # =========================================================================
    add_section_slide(prs, "Supplementary Material", "S")
    
    # =========================================================================
    # SLIDE S1: SAE Equations
    # =========================================================================
    add_content_slide(prs, "SAE Mathematical Formulation", [
        "Encoder: h = ReLU(W_enc * x + b_enc)",
        "   -> Maps d_model -> n_features (8x expansion)",
        "   -> ReLU enforces sparsity",
        "Decoder: x_hat = W_dec * h + b_dec",
        "   -> Reconstructs original activations",
        "Loss: L = MSE(x, x_hat) + lambda*||h||_1",
        "   -> Reconstruction fidelity + sparsity penalty",
        "   -> lambda = 10^-4 (PaLiGemma), 5x10^-4 (Qwen)",
        "Feature identification: |mean(h_j^male) - mean(h_j^female)|",
        "   -> Top-k features selected (k=100)"
    ])
    
    # =========================================================================
    # SLIDE S2: Ablation Equations
    # =========================================================================
    add_content_slide(prs, "Ablation Hook Methods", [
        "FULL RECONSTRUCTION (PaLiGemma):",
        "   x_mod = W_dec * (h element-wise m) + b_dec",
        "   -> m is binary mask (0 for ablated features)",
        "   -> Replaces activations with SAE output",
        "   -> Works when SAE has high reconstruction quality",
        "RESIDUAL ABLATION (Qwen, Llama):",
        "   x_mod = x - Sum_{j in A} h_j * W_dec^(:,j)",
        "   -> A is set of ablated feature indices",
        "   -> Subtracts only targeted feature contribution",
        "   -> Preserves original activations otherwise",
        "   -> Needed when SAE trained on pooled but applied to per-token"
    ])
    
    # =========================================================================
    # SLIDE S3: Statistical Methods
    # =========================================================================
    add_content_slide(prs, "Statistical Methods", [
        "Primary metric: Total gender terms across 500 captions",
        "Percentage change from baseline reported",
        "PAIRED ANALYSIS:",
        "   -> Per-image: Delta_i = count(ablated_i) - count(baseline_i)",
        "   -> Bootstrap 95% CI (10,000 resamples)",
        "   -> Wilcoxon signed-rank test",
        "EFFECT SPECIFICITY:",
        "   -> Specificity = |targeted%| - |random%|",
        "   -> Ratio = |targeted| / |random|",
        "   -> Values > 1 indicate targeted effect exceeds random",
        "Significance: CI excluding zero confirms effect"
    ])
    
    # =========================================================================
    # SLIDE S4: LLaVA Exclusion
    # =========================================================================
    add_content_slide(prs, "Why LLaVA Was Not Used for Intervention", [
        "LLaVA-1.5-7B trained but excluded from intervention experiments",
        "Reason: 95% dead SAE features",
        "   -> Most features never activate",
        "   -> Insufficient active features for meaningful ablation",
        "Secondary reason: Architectural similarity to Qwen2-VL",
        "   -> Both are instruction-tuned 7B models",
        "   -> Would likely show same direction (inhibitory)",
        "   -> Qwen provides sufficient evidence for this class",
        "LLaVA results included in cross-lingual feature overlap analysis",
        "   -> Confirms near-zero overlap pattern (Jaccard = 0.010)"
    ])
    
    # =========================================================================
    # SLIDE S5: Arabic Per-Term
    # =========================================================================
    add_table_slide(prs, "Arabic Per-Term Changes (AR->AR Ablation)",
        ["Term (Transliteration)", "Meaning", "PaLiGemma", "Qwen2-VL"],
        [
            ["rajul", "man", "+5.0%", "-57.4%"],
            ["imra'a", "woman", "-13.3%", "-24.2%"],
            ["huwa", "he", "+14.3%", "+0.0%"],
            ["hiya", "she", "-33.3%", "-27.3%"],
            ["ab", "father", "-42.1%", "+12.0%"],
            ["umm", "mother", "-9.5%", "+20.6%"],
            ["akh", "brother", "+0.0%", "+32.4%"]
        ]
    )
    
    # =========================================================================
    # SLIDE S6: Related Work
    # =========================================================================
    add_content_slide(prs, "Related Work", [
        "GENDER BIAS IN VLMs:",
        "   -> Hendricks et al. 2018: Activity vs appearance bias",
        "   -> Zhao et al. 2017: Data augmentation for debiasing",
        "   -> Hirota et al. 2022: Bias quantification metrics",
        "MECHANISTIC INTERPRETABILITY:",
        "   -> Cunningham et al. 2023: SAEs on GPT-2",
        "   -> Bricken et al. 2023: Monosemanticity at scale",
        "   -> Templeton et al. 2024: 34M features on Claude",
        "SAE ARCHITECTURES:",
        "   -> Gao et al. 2024: TopK SAEs",
        "   -> Rajamanoharan et al. 2024: JumpReLU SAEs",
        "   -> Karvonen et al. 2025: SAEBench evaluation"
    ])
    
    # Save presentation
    output_path = "/home2/jmsk62/mechanistic_intrep/sae_captioning_project/presentation/SAE_Gender_Bias_VLM_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    print(f"   Total slides: {len(prs.slides)}")
    
    return output_path

if __name__ == "__main__":
    create_presentation()
